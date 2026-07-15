from asyncio import Lock
from dataclasses import dataclass, field

from pptx.presentation import Presentation as PresentationType
from pydantic import BaseModel, Field


@dataclass
class Project:
    presentation: PresentationType
    master_name: str | None = None
    lock: Lock = field(default_factory=Lock)


class SlideInfo(BaseModel):
    layout: str = Field(
        description="Layout name the slide is based on.",
    )
    text: str = Field(
        description="All text currently on the slide.",
    )


class TemplatesResult(BaseModel):
    hint: str = Field(
        description=(
            "Suggested next step — guidance for the agent, not part of the data."
        ),
    )
    templates: list[str] = Field(
        description="Template file names for `create_project`.",
    )


class ProjectResult(BaseModel):
    hint: str = Field(
        description=(
            "Suggested next step — guidance for the agent, not part of the data."
        ),
    )
    slide_count: int = Field(
        description="Current number of slides in the project.",
    )


class StartResult(ProjectResult):
    masters: list[str] = Field(
        description=(
            "Master names. Pick one (ask the user if unclear), inspect it "
            "with `list_layouts`, then select it via `set_master` in "
            "`run_script`."
        ),
    )


class LayoutsResult(BaseModel):
    hint: str = Field(
        description=(
            "Suggested next step — guidance for the agent, not part of the data."
        ),
    )
    layouts: dict[str, dict[int, str]] = Field(
        description=(
            "Layout name (for `add_slide`) -> placeholder idx -> placeholder "
            "type, e.g. TITLE or BODY. Target placeholders by idx in `fill`, "
            "`add_image` and `add_chart`."
        ),
    )


class ScriptResult(ProjectResult):
    output: str = Field(
        description="Printed output and last expression of the script.",
    )


class FinalizeResult(ProjectResult):
    file_name: str = Field(
        description="Name of the uploaded file.",
    )
    owui_url: str = Field(
        description="OpenWebUI download URL of the uploaded file.",
    )
