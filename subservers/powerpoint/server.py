import asyncio
import io
from collections.abc import AsyncIterator
from contextlib import asynccontextmanager, suppress
from dataclasses import dataclass, field

import httpx
from cachetools import TTLCache
from fastmcp import FastMCP
from fastmcp.server.dependencies import get_access_token
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pydantic import Field

from config import get_settings
from models.project import ProjectInfo, SavedProjectInfo, SlideInfo
from models.template import TemplateInfo
from services.owui import upload_file
from subservers.powerpoint._utils import (
    analyze_templates,
    drop_all_slides,
    drop_slide,
)

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

@dataclass
class _Project:
    template_name: str
    pptx: PresentationType
    lock: asyncio.Lock = field(default_factory=asyncio.Lock)


_settings = get_settings()
_projects: TTLCache[str, _Project] = TTLCache(
    maxsize=10_000, ttl=_settings.project_ttl_seconds
)
_templates: dict[str, TemplateInfo] = {}


async def _sweep_projects(interval: float) -> None:
    while True:
        await asyncio.sleep(interval)
        _projects.expire()


@asynccontextmanager
async def lifespan(server: FastMCP) -> AsyncIterator[None]:
    global _templates
    _templates = analyze_templates(_settings.templates_dir)
    sweeper = asyncio.create_task(
        _sweep_projects(_settings.project_sweep_interval_seconds)
    )
    try:
        yield
    finally:
        sweeper.cancel()
        with suppress(asyncio.CancelledError):
            await sweeper


mcp = FastMCP(name="powerpoint", lifespan=lifespan)


@mcp.tool(
    name="list_templates",
    description=(
        "List all PowerPoint templates (.pptx) that were loaded and analyzed from "
        "the configured templates directory at server startup. Returns a mapping "
        "of template name to its info (absolute path, slide count, and all slide "
        "layouts including their placeholders: index, name, type). Use this tool "
        "to discover which templates are available and which layouts/placeholders "
        "can be filled when generating new slides."
    ),
)
async def list_templates() -> dict[str, TemplateInfo]:
    return _templates


@mcp.tool(
    name="create_project",
    description=(
        "Create a new in-memory PowerPoint project for the calling user, based "
        "on one of the available templates. The template is loaded fresh into "
        "memory and stored under the user's id (from the JWT), so subsequent "
        "tool calls by the same user can modify it. Pass the template `name` "
        "exactly as returned by `list_templates`. If a project already exists "
        "for this user it is overwritten."
    ),
)
async def create_project(
    template_name: str = Field(
        description="Name of the template to use, as returned by `list_templates`."
    ),
) -> ProjectInfo:
    template = _templates.get(template_name)
    if template is None:
        raise ValueError(
            f"Template '{template_name}' not found."
            f" Available: {list(_templates)}"
        )

    token = get_access_token()
    if token is None:
        raise ValueError("Authentication required.")
    user_id = token.claims.get("id")
    if not user_id:
        raise ValueError("JWT is missing the 'id' claim.")
    user_id = str(user_id)

    pptx = await asyncio.to_thread(Presentation, template.path)
    drop_all_slides(pptx)
    _projects[user_id] = _Project(template_name=template_name, pptx=pptx)

    return ProjectInfo(
        user_id=user_id,
        template_name=template_name,
        slide_count=len(pptx.slides),
    )


@mcp.tool(
    name="append_slide",
    description=(
        "Append a new slide to the calling user's project. The slide is created "
        "from one of the slide layouts defined in the project's template (see "
        "`list_templates` for layout names and their placeholders). Optionally "
        "fill text placeholders by passing a mapping from placeholder index "
        "(`idx`) to the text value to insert; placeholders not in the mapping "
        "are left empty. Returns the zero-based index of the appended slide "
        "within the project."
    ),
)
async def append_slide(
    layout_name: str = Field(
        description="Name of the slide layout to use for the new slide."
    ),
    placeholders: dict[int, str] = Field(
        default_factory=dict,
        description=(
            "Optional mapping from placeholder index (`idx`) to the text value to insert."
        ),
    ),
) -> SlideInfo:
    token = get_access_token()
    if token is None:
        raise ValueError("Authentication required.")
    user_id = token.claims.get("id")
    if not user_id:
        raise ValueError("JWT is missing the 'id' claim.")
    user_id = str(user_id)

    project = _projects.get(user_id)
    if project is None:
        raise ValueError(
            "No project for this user."
            " Call `create_project` first."
        )

    async with project.lock:
        layout = project.pptx.slide_layouts.get_by_name(layout_name)
        if layout is None:
            available = [lo.name for lo in project.pptx.slide_layouts]
            raise ValueError(
                f"Layout '{layout_name}' not found in template "
                f"'{project.template_name}'. Available: {available}"
            )

        slide = project.pptx.slides.add_slide(layout)

        for idx, text in placeholders.items():
            with suppress(KeyError):
                slide.placeholders[idx].text = text

        _projects[user_id] = project

        return SlideInfo(
            index=len(project.pptx.slides) - 1,
            layout_name=layout_name,
        )


@mcp.tool(
    name="edit_slide",
    description=(
        "Update text placeholders on an existing slide of the calling user's "
        "project, identified by zero-based index. Pass a mapping from "
        "placeholder index (`idx`) to the new text value; only those "
        "placeholders are modified, the rest are left as-is. Returns the "
        "slide's index and its layout name."
    ),
)
async def edit_slide(
    index: int = Field(
        description="Zero-based index of the slide to edit."
    ),
    placeholders: dict[int, str] = Field(
        description=(
            "Mapping from placeholder index (`idx`) to the new text value. "
            "Only keys listed here are updated; other placeholders stay as-is."
        ),
    ),
) -> SlideInfo:
    token = get_access_token()
    if token is None:
        raise ValueError("Authentication required.")
    user_id = token.claims.get("id")
    if not user_id:
        raise ValueError("JWT is missing the 'id' claim.")
    user_id = str(user_id)

    project = _projects.get(user_id)
    if project is None:
        raise ValueError(
            "No project for this user."
            " Call `create_project` first."
        )

    async with project.lock:
        count = len(project.pptx.slides)
        if not 0 <= index < count:
            raise ValueError(
                f"Slide index {index} out of range."
                f" Project has {count} slide(s) (valid: 0..{count - 1})."
            )

        slide = project.pptx.slides[index]
        for idx, text in placeholders.items():
            with suppress(KeyError):
                slide.placeholders[idx].text = text

        _projects[user_id] = project

        return SlideInfo(
            index=index,
            layout_name=slide.slide_layout.name,
        )


@mcp.tool(
    name="remove_slides",
    description=(
        "Remove one or more slides from the calling user's project by their "
        "zero-based indices. Indices refer to the slide positions before "
        "removal; the call validates all indices first, then removes them in "
        "descending order so earlier indices stay stable. Returns the updated "
        "project info with the new slide count."
    ),
)
async def remove_slides(
    indices: list[int] = Field(
        description=(
            "Zero-based slide indices to remove. Duplicates are ignored."
        ),
    ),
) -> ProjectInfo:
    token = get_access_token()
    if token is None:
        raise ValueError("Authentication required.")
    user_id = token.claims.get("id")
    if not user_id:
        raise ValueError("JWT is missing the 'id' claim.")
    user_id = str(user_id)

    project = _projects.get(user_id)
    if project is None:
        raise ValueError(
            "No project for this user."
            " Call `create_project` first."
        )

    async with project.lock:
        count = len(project.pptx.slides)
        out_of_range = [i for i in indices if not 0 <= i < count]
        if out_of_range:
            raise ValueError(
                f"Slide indices out of range {out_of_range}."
                f" Project has {count} slide(s) (valid: 0..{count - 1})."
            )

        for i in sorted(set(indices), reverse=True):
            drop_slide(project.pptx, i)

        _projects[user_id] = project

        return ProjectInfo(
            user_id=user_id,
            template_name=project.template_name,
            slide_count=len(project.pptx.slides),
        )


@mcp.tool(
    name="save_project",
    description=(
        "Serialize the calling user's project to a `.pptx` byte stream and "
        "upload it to OpenWebUI using the caller's JWT. Pass `filename` "
        "without extension (`.pptx` is appended automatically); if omitted, "
        "the template name is used. The project stays active after saving, "
        "so further `append_slide` calls and re-saving are possible. Returns "
        "the chosen filename, the slide count, the OpenWebUI file id, and a "
        "download URL."
    ),
)
async def save_project(
    filename: str = Field(
        default="",
        description=(
            "File name without extension. Defaults to the template name "
            "when empty. Existing files with the same name are overwritten."
        ),
    ),
) -> SavedProjectInfo:
    token = get_access_token()
    if token is None:
        raise ValueError("Authentication required.")
    user_id = token.claims.get("id")
    if not user_id:
        raise ValueError("JWT is missing the 'id' claim.")
    user_id = str(user_id)

    project = _projects.get(user_id)
    if project is None:
        raise ValueError(
            "No project for this user."
            " Call `create_project` first."
        )

    async with project.lock:
        stem = filename.strip() or project.template_name
        out_name = f"{stem}.pptx"

        buf = io.BytesIO()
        await asyncio.to_thread(project.pptx.save, buf)

        _projects[user_id] = project
        slide_count = len(project.pptx.slides)

    base_url = _settings.owui_base_url.rstrip("/")
    try:
        uploaded = await upload_file(
            filename=out_name,
            data=buf.getvalue(),
            content_type=PPTX_MIME,
            token=token.token,
            base_url=base_url,
        )
    except httpx.HTTPStatusError as exc:
        raise RuntimeError(
            f"OpenWebUI rejected the upload (HTTP {exc.response.status_code})."
            f" Response body: {exc.response.text[:500]}"
        ) from exc
    except httpx.RequestError as exc:
        raise RuntimeError(
            f"Could not reach OpenWebUI at {base_url}."
            f" Check OWUI_BASE_URL and that the service is running. Detail: {exc}"
        ) from exc

    return SavedProjectInfo(
        filename=out_name,
        slide_count=slide_count,
        owui_file_id=uploaded.id,
        owui_url=f"{base_url}/api/v1/files/{uploaded.id}/content",
    )
