from fastmcp.utilities.logging import get_logger
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.oxml import parse_xml
from pptx.presentation import Presentation as PresentationType

from models.template import LayoutInfo, PlaceholderInfo


logger = get_logger(__name__)


def list_template_names(
    templates_dir: Path,
) -> list[str]:

    if not templates_dir.is_dir():
        raise RuntimeError(f"Templates directory not found: {templates_dir}")

    template_names: list[str] = []

    for file_path in sorted(templates_dir.glob("*.pptx")):
        try:
            _ = Presentation(file_path)
        except Exception as error:
            logger.warning(f"Skipping template {file_path.name}: {error}")
        else:
            template_names.append(file_path.name)

    return template_names


def list_master_names(
    templates_dir: Path,
    template_name: str,
) -> dict[int, str]:

    presentation = Presentation(templates_dir.joinpath(template_name))

    master_names: dict[int, str] = {}

    for i, master in enumerate(presentation.slide_masters):

        master_name = master.name

        if not master_name:
            theme_part = master.part.part_related_by(RELATIONSHIP_TYPE.THEME)
            theme_xml = parse_xml(theme_part.blob)
            master_name = theme_xml.get("name")

        if master_name:
            master_names[i] = master_name

    return master_names


def list_layout_infos(
    templates_dir: Path,
    template_name: str,
    master_index: int,
) -> dict[str, LayoutInfo]:
    
    presentation = Presentation(templates_dir.joinpath(template_name))

    master = presentation.slide_masters[master_index]

    layout_infos: dict[str, LayoutInfo] = {}

    for layout in master.slide_layouts:

        placeholders: list[PlaceholderInfo] = []

        for placeholder in layout.placeholders:

            placeholders.append(
                PlaceholderInfo(
                    idx=placeholder.placeholder_format.idx,
                    name=placeholder.name,
                    type=str(placeholder.placeholder_format.type),
                )
            )

        layout_infos[layout.name] = LayoutInfo(placeholders=placeholders)

    return layout_infos


_RID_ATTR = (
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
)


def drop_slide(presentation: PresentationType, index: int) -> None:

    slide_id_lst = presentation.slides._sldIdLst

    try:
        slide_id = list(slide_id_lst)[index]
    except IndexError:
        raise ValueError(
            f"Slide index {index} out of range."
        )

    if rid := slide_id.get(_RID_ATTR):
        presentation.part.drop_rel(rid)

    slide_id_lst.remove(slide_id)


def drop_all_slides(presentation: PresentationType) -> None:
    while len(presentation.slides._sldIdLst) > 0:
        drop_slide(presentation, 0)
